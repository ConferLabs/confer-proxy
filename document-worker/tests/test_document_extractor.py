from __future__ import annotations

import io
import sys

import pytest
from docx import Document
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from confer_document_worker.artifact.ArtifactReader import ArtifactReader
from confer_document_worker.artifact.ArtifactWriter import ArtifactWriter
from confer_document_worker.domain.ContainerKind import ContainerKind
from confer_document_worker.domain.SearchQuery import SearchQuery
from confer_document_worker.extraction.DocumentExtractorRouter import DocumentExtractorRouter
from confer_document_worker.extraction.MediaTypes import (
  DOCX_MEDIA_TYPE,
  PPTX_MEDIA_TYPE,
  XLSX_MEDIA_TYPE,
)
from confer_document_worker.extraction.office.DoclingDocumentConverter import DoclingDocumentConverter
from confer_document_worker.extraction.office.DoclingDocumentMapper import DoclingDocumentMapper
from confer_document_worker.extraction.office.DoclingSlimDocumentExtractor import DoclingSlimDocumentExtractor
from confer_document_worker.extraction.office.OfficeArchivePreflight import OfficeArchivePreflight
from confer_document_worker.extraction.office.OfficeExtractionLimits import OfficeExtractionLimits
from confer_document_worker.extraction.text.TextDocumentExtractor import TextDocumentExtractor
from SourceFixture import SourceFixture
from test_extractor import pdf_extractor


def document_extractor_router() -> DocumentExtractorRouter:
  limits = OfficeExtractionLimits()
  return DocumentExtractorRouter(
    pdf_extractor(),
    DoclingSlimDocumentExtractor(
      OfficeArchivePreflight(),
      DoclingDocumentConverter(limits),
      DoclingDocumentMapper(limits),
      limits,
    ),
    TextDocumentExtractor(),
  )


def extract(content: bytes, filename: str, media_type: str):
  # Protocol payloads deliberately use opaque .bin spool names. Extractors must
  # route by the authenticated request metadata, never by a temporary suffix.
  with SourceFixture(content, suffix=".bin") as source:
    return document_extractor_router().extract(source, filename, media_type)


def docx_document() -> bytes:
  output = io.BytesIO()
  document = Document()
  document.add_heading("Health history", level=1)
  document.add_paragraph("Analyte increased between January and June.")
  table = document.add_table(rows=3, cols=2)
  table.cell(0, 0).text = "Date"
  table.cell(0, 1).text = "Result"
  table.cell(1, 0).text = "January"
  table.cell(1, 1).text = "410 ng/dL"
  table.cell(2, 0).text = "June"
  table.cell(2, 1).text = "525 ng/dL"
  document.save(output)
  return output.getvalue()


def pptx_document() -> bytes:
  output = io.BytesIO()
  image_output = io.BytesIO()
  Image.new("RGB", (32, 24), (20, 90, 160)).save(image_output, format="PNG")
  image_output.seek(0)

  presentation = Presentation()
  slide = presentation.slides.add_slide(presentation.slide_layouts[5])
  textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
  textbox.text = "Quarterly results"
  table = slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(5), Inches(2)).table
  table.cell(0, 0).text = "Quarter"
  table.cell(0, 1).text = "Revenue"
  table.cell(1, 0).text = "Q1"
  table.cell(1, 1).text = "$125"
  slide.shapes.add_picture(image_output, Inches(6), Inches(1), width=Inches(1))
  presentation.save(output)
  return output.getvalue()


def xlsx_document() -> bytes:
  output = io.BytesIO()
  workbook = Workbook()
  labs = workbook.active
  labs.title = "Labs"
  labs.append(["Test", "Result"])
  labs.append(["Vitamin D", 32])
  history = workbook.create_sheet("History")
  history.append(["Date", "Analyte"])
  history.append(["June", 525])
  workbook.save(output)
  return output.getvalue()


def test_docx_uses_native_headings_and_tables() -> None:
  content = extract(docx_document(), "history.docx", DOCX_MEDIA_TYPE)

  assert content.media_type == DOCX_MEDIA_TYPE
  assert content.containers[0].kind == ContainerKind.DOCUMENT
  assert content.outline[0].title == "Health history"
  assert "525 ng/dL" in "\n".join(block.text for block in content.containers[0].blocks)
  assert content.containers[0].regions[0].type.value == "table"
  assert not content.containers[0].regions[0].renderable


def test_pptx_preserves_slide_coordinates_and_embedded_images() -> None:
  source = pptx_document()
  content = extract(source, "results.pptx", PPTX_MEDIA_TYPE)
  picture = next(region for region in content.containers[0].regions if region.type.value == "image")

  assert content.containers[0].kind == ContainerKind.SLIDE
  assert "Quarterly results" in "\n".join(block.text for block in content.containers[0].blocks)
  assert picture.renderable
  assert picture.asset is not None and picture.asset.startswith(b"\x89PNG")



def test_xlsx_preserves_sheet_names_and_native_cell_tables() -> None:
  content = extract(xlsx_document(), "records.xlsx", XLSX_MEDIA_TYPE)

  assert [page.label for page in content.containers] == ["Labs", "History"]
  assert all(page.kind == ContainerKind.SHEET for page in content.containers)
  assert "Vitamin D" in content.containers[0].blocks[0].text
  assert "Analyte" in content.containers[1].blocks[0].text


def test_text_and_csv_receive_searchable_artifacts() -> None:
  text = extract(b"First note\n\nAnalyte 525 ng/dL", "notes.txt", "text/plain")
  csv = extract(b"Test,Result\nVitamin D,32\n", "labs.csv", "text/csv")

  with ArtifactReader(ArtifactWriter().build(text)) as reader:
    assert reader.search((SearchQuery(("analyte",)),))[0].text.endswith("525 ng/dL")
  with ArtifactReader(ArtifactWriter().build(csv)) as reader:
    assert reader.search((SearchQuery(("Vitamin D",)),))[0].regions[0].type.value == "table"
    assert not reader.search((SearchQuery(("Vitamin D",)),))[0].regions[0].renderable


def test_html_and_markdown_preserve_headings() -> None:
  html = extract(b"<h1>Summary</h1><p>Native HTML text</p>", "page.html", "text/html")
  markdown = extract(b"# Summary\n\nNative Markdown text", "notes.md", "text/markdown")

  assert html.outline[0].title == "Summary"
  assert markdown.outline[0].title == "Summary"


def test_office_preflight_rejects_invalid_archives() -> None:
  with pytest.raises(ValueError, match="archive is invalid"):
    extract(b"not a zip", "record.docx", DOCX_MEDIA_TYPE)


def test_runtime_check_rejects_a_missing_structured_document_backend(
  monkeypatch: pytest.MonkeyPatch,
) -> None:
  monkeypatch.setitem(sys.modules, "docling.backend.msexcel_backend", None)

  with pytest.raises(ModuleNotFoundError):
    DoclingDocumentConverter(OfficeExtractionLimits()).check()
